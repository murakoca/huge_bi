from pptx import Presentation
from pptx.util import Inches, Pt
import plotly.io as pio
import pandas as pd
import io
from ai.report_writer import ReportWriter  # Part 5’teki LLM özetleyici

class StoryTeller:
    def __init__(self, llm_api_key=None):
        self.writer = ReportWriter(api_key=llm_api_key)

    def create_story(self, data_frames: dict, output_path="story.pptx"):
        """
        data_frames: {'slayt_başlığı': (DataFrame, grafik_türü)}
        Önce LLM ile metin üretir, sonra slayt destesini oluşturur.
        """
        prs = Presentation()
        for title, (df, chart_type) in data_frames.items():
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Boş slayt
            # Başlık
            slide.shapes.title.text = title
            # LLM ile özet paragraf
            summary = self.writer.summarize(df)
            txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2))
            tf = txBox.text_frame
            tf.text = summary
            # Grafik (plotly image)
            fig = self._get_figure(df, chart_type)
            img_bytes = pio.to_image(fig, format='png')
            img_stream = io.BytesIO(img_bytes)
            slide.shapes.add_picture(img_stream, Inches(1), Inches(3.5), Inches(8), Inches(4.5))
        prs.save(output_path)
        return output_path

    def _get_figure(self, df, chart_type):
        import plotly.express as px
        if chart_type == 'bar':
            fig = px.bar(df, x=df.columns[0], y=df.columns[1])
        elif chart_type == 'line':
            fig = px.line(df, x=df.columns[0], y=df.columns[1])
        elif chart_type == 'pie':
            fig = px.pie(df, names=df.columns[0], values=df.columns[1])
        else:
            fig = px.scatter(df, x=df.columns[0], y=df.columns[1])
        fig.update_layout(margin=dict(l=0,r=0,t=0,b=0))
        return fig